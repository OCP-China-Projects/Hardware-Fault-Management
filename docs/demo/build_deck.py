#!/usr/bin/env python3
"""Generate the Hardware Fault Management RAS API demo deck (.pptx).

Content is grounded in the RASAPI repo (README.md, docs/pldm-mctp-i3c-design.md)
and verified live on the internal dev machine. No placeholder data.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ---------------------------------------------------------------
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)   # deep background
STEEL  = RGBColor(0x1E, 0x3A, 0x5F)   # panel
ACCENT = RGBColor(0x35, 0xC2, 0xD1)   # cyan accent
AMBER  = RGBColor(0xF2, 0xA6, 0x3B)   # highlight
GREEN  = RGBColor(0x4C, 0xC9, 0x6A)   # pass
RED    = RGBColor(0xE5, 0x5A, 0x5A)   # caution
WHITE  = RGBColor(0xF4, 0xF7, 0xFA)
GREY   = RGBColor(0xB8, 0xC4, 0xD0)
CODEBG = RGBColor(0x10, 0x18, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, italic, font)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = font
    return tb


def R(txt, size=18, color=WHITE, bold=False, italic=False, font="Calibri"):
    return (txt, size, color, bold, italic, font)


def C(txt, size=13, color=ACCENT):
    return (txt, size, color, False, False, "Consolas")


def header(slide, kicker, title):
    box(slide, 0, 0, 13.333, 1.35, fill=STEEL)
    box(slide, 0, 1.35, 13.333, 0.05, fill=ACCENT)
    text(slide, 0.6, 0.18, 12, 0.4, [[R(kicker, 13, ACCENT, True)]])
    text(slide, 0.6, 0.5, 12.1, 0.8, [[R(title, 27, WHITE, True)]])


def footer(slide, n):
    text(slide, 0.6, 7.05, 8, 0.35,
         [[R("Hardware Fault Management · RAS API · MCTP/PLDM over QEMU", 10, GREY)]])
    text(slide, 12.2, 7.05, 0.8, 0.35, [[R(str(n), 10, GREY)]], align=PP_ALIGN.RIGHT)


# ============================================================ Slide 1 — Title
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
box(s, 0, 2.55, 13.333, 0.06, fill=ACCENT)
box(s, 0.9, 2.75, 6.2, 0.06, fill=AMBER)
text(s, 0.9, 0.9, 11.5, 0.5, [[R("OCP CHINA · HARDWARE FAULT MANAGEMENT", 16, ACCENT, True)]])
text(s, 0.9, 1.4, 11.9, 1.2,
     [[R("Hardware Fault Management RAS API", 40, WHITE, True)]])
text(s, 0.9, 2.75, 11.9, 0.9,
     [[R("PLDM over MCTP end-to-end demo — OpenBMC (AST2600) \u2194 Zephyr (RISC-V), bridged over two live QEMU instances", 18, GREY)]])
# quick fact strip
facts = [
    ("Requester", "OpenBMC pldmd / mctpd — EID 8"),
    ("Responder", "Zephyr serial_bridge — EID 18"),
    ("Transport", "MCTP-over-serial (DSP0253) via unix socket"),
    ("Result", "Auto-discovery + fwd/reverse PLDM (rc=0)"),
]
x = 0.9
for lab, val in facts:
    box(s, x, 3.95, 2.85, 1.5, fill=STEEL, line=ACCENT, line_w=0.75)
    text(s, x+0.2, 4.1, 2.5, 0.4, [[R(lab, 13, ACCENT, True)]])
    text(s, x+0.2, 4.5, 2.5, 0.9, [[R(val, 13.5, WHITE)]], line_spacing=1.05)
    x += 3.03
text(s, 0.9, 5.9, 11.9, 0.6,
     [[R("Demo host: internal dev machine   ·   branch RASAPI", 14, GREY)]])
text(s, 0.9, 6.5, 11.9, 0.5,
     [[R("Presenter guide: see DEMO_GUIDE.md for the exact step-by-step SSH walkthrough", 12, ACCENT, italic=True)]])

# ============================================================ Slide 2 — Agenda
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "OVERVIEW", "What we will demonstrate")
items = [
    ("1", "The problem & the RAS API goal", "Why HFM needs a standard sensor/effecter transport between a BMC and satellite MCUs"),
    ("2", "Architecture — the four moving pieces", "QEMU 11 · Zephyr on RISC-V · I3C model · MCTP+PLDM stack"),
    ("3", "The realized topology", "Two QEMU instances bridged over MCTP-serial (DSP0253)"),
    ("4", "Live demo — forward discovery", "OpenBMC pldmtool drives PLDM Type 0/2 against Zephyr EID 18"),
    ("5", "Live demo — reverse path (patch 0009)", "Zephyr requester \u2192 BMC responder, tag-owner bit fix"),
    ("6", "Full auto-discovery — three fixes & wrap-up", "kernel busy-loop (0010) + control TO bit (0011) + systemd order (0006)"),
]
y = 1.75
for num, title, sub in items:
    box(s, 0.7, y, 0.7, 0.7, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    text(s, 0.7, y+0.06, 0.7, 0.6, [[R(num, 22, NAVY, True)]], align=PP_ALIGN.CENTER)
    text(s, 1.65, y-0.02, 11, 0.45, [[R(title, 19, WHITE, True)]])
    text(s, 1.65, y+0.4, 11, 0.4, [[R(sub, 13.5, GREY)]])
    y += 0.87
footer(s, 2)

# ============================================================ Slide 3 — Problem/Goal
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "MOTIVATION", "Why a RAS API — and why PLDM over MCTP")
box(s, 0.7, 1.7, 5.9, 4.9, fill=STEEL)
text(s, 0.95, 1.85, 5.4, 0.4, [[R("The pain point", 16, AMBER, True)]])
text(s, 0.95, 2.35, 5.45, 4.1, [
    [R("Large-scale data centers must detect, report and act on hardware faults across thousands of heterogeneous devices.", 14.5, WHITE)],
    [R("", 6, WHITE)],
    [R("A BMC needs a vendor-neutral way to:", 14.5, GREY)],
    [R("•  discover satellite controllers", 14, WHITE)],
    [R("•  read sensors (temp, voltage, fan)", 14, WHITE)],
    [R("•  read/set thresholds & effecters", 14, WHITE)],
    [R("•  receive async fault events", 14, WHITE)],
], line_spacing=1.12)
box(s, 6.75, 1.7, 5.9, 4.9, fill=STEEL)
text(s, 7.0, 1.85, 5.4, 0.4, [[R("The DMTF PMCI answer", 16, ACCENT, True)]])
text(s, 7.0, 2.35, 5.45, 4.1, [
    [R("PLDM (Platform Level Data Model) carried over MCTP (Management Component Transport Protocol).", 14.5, WHITE)],
    [R("", 6, WHITE)],
    [R("This demo proves the full stack on emulated hardware:", 14.5, GREY)],
    [R("PLDM Type 0", 14, ACCENT, True), R("  — base / discovery (DSP0240)", 14, WHITE)],
    [R("PLDM Type 2", 14, ACCENT, True), R("  — platform monitoring, PDR & sensors (DSP0248)", 14, WHITE)],
    [R("MCTP", 14, ACCENT, True), R("  — transport & EID routing (DSP0236)", 14, WHITE)],
], line_spacing=1.12)
footer(s, 3)

# ============================================================ Slide 4 — Layer stack
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "ARCHITECTURE", "The protocol layer stack")
layers = [
    ("App — PLDM Type 0/2 responder  (samples/.../serial_bridge)", ACCENT),
    ("libpldm — encode/decode  (vendored openbmc/libpldm)", STEEL),
    ("libmctp — framing, EID routing  (upstream module)", STEEL),
    ("MCTP bindings — serial DSP0253 · I3C DSP0233 · I2C DSP0237", STEEL),
    ("Zephyr UART / I3C driver  (upstream)", STEEL),
    ("QEMU 11 machine model  (sifive_u / ast2600-evb)", STEEL),
]
y = 1.85
for i, (lab, col) in enumerate(layers):
    fill = ACCENT if i == 0 else STEEL
    txtcol = NAVY if i == 0 else WHITE
    box(s, 2.3, y, 8.7, 0.72, fill=fill, line=ACCENT, line_w=0.75)
    text(s, 2.55, y+0.12, 8.3, 0.5, [[R(lab, 15, txtcol, i == 0)]], anchor=MSO_ANCHOR.MIDDLE)
    if i < len(layers)-1:
        text(s, 6.4, y+0.66, 0.6, 0.3, [[R("\u25bc", 12, GREY)]], align=PP_ALIGN.CENTER)
    y += 0.82
text(s, 0.7, 6.7, 12, 0.4,
     [[R("Delivered as numbered patches 0001–0011 under patches/ — full apply order in patches/README.md", 13, GREY, italic=True)]])
footer(s, 4)

# ============================================================ Slide 5 — Topology diagram
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "ARCHITECTURE", "Realized topology — two QEMU instances over MCTP-serial")
# left box: OpenBMC
box(s, 0.7, 2.0, 5.2, 3.5, fill=STEEL, line=AMBER, line_w=1.5)
text(s, 0.9, 2.15, 4.8, 0.5, [[R("OpenBMC QEMU", 18, AMBER, True)], [R("-machine ast2600-evb", 12, GREY, False, False, "Consolas")]])
for i, ln in enumerate(["pldmd  (PLDM Type 0/2 requester)",
                        "mctpd  (AF_MCTP bus-owner)",
                        "kernel mctp-serial driver",
                        "EID 8"]):
    box(s, 0.95, 3.05+i*0.55, 4.7, 0.45, fill=CODEBG, line=STEEL, line_w=0.5)
    col = ACCENT if "EID" in ln else WHITE
    text(s, 1.1, 3.08+i*0.55, 4.5, 0.4, [[R(ln, 13.5, col, "EID" in ln)]], anchor=MSO_ANCHOR.MIDDLE)
# right box: Zephyr
box(s, 7.45, 2.0, 5.2, 3.5, fill=STEEL, line=ACCENT, line_w=1.5)
text(s, 7.65, 2.15, 4.8, 0.5, [[R("Zephyr QEMU", 18, ACCENT, True)], [R("-machine sifive_u", 12, GREY, False, False, "Consolas")]])
for i, ln in enumerate(["serial_bridge sample",
                        "PLDM Type 0 + Type 2 responder",
                        "MCTP control responder",
                        "EID 18"]):
    box(s, 7.7, 3.05+i*0.55, 4.7, 0.45, fill=CODEBG, line=STEEL, line_w=0.5)
    col = ACCENT if "EID" in ln else WHITE
    text(s, 7.85, 3.08+i*0.55, 4.5, 0.4, [[R(ln, 13.5, col, "EID" in ln)]], anchor=MSO_ANCHOR.MIDDLE)
# link arrow
box(s, 5.9, 3.55, 1.55, 0.45, fill=GREEN, line=None, shape=MSO_SHAPE.LEFT_RIGHT_ARROW)
text(s, 5.55, 5.65, 6.2, 0.9,
     [[R("unix socket  /tmp/hfm-mctp.sock", 14, GREEN, True, False, "Consolas")],
      [R("OpenBMC = socket server (start first) · Zephyr = client", 12.5, GREY)]],
     align=PP_ALIGN.CENTER)
text(s, 0.7, 6.55, 12, 0.5,
     [[R("MCTP link rides each side's SECOND -serial backend (uart1 / ttyS0); the first -serial stays the console.", 13, GREY, italic=True)]])
footer(s, 5)

# ============================================================ Slide 6 — Build artifacts
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "PREREQUISITES", "What is already built on the dev machine")
rows = [
    ("Component", "Location / artifact", "Verify", True),
    ("QEMU 11.0.0", "~/qemu-build/bin/qemu-system-{arm,riscv64}", "runs -machine help", False),
    ("OpenBMC image (+0004/6/10)", "obmc-phosphor-image-evb-ast2600.static.mtd", "md5 9a858c13", False),
    ("Zephyr ELF (+0009/0011)", "build-serial-fix/zephyr/zephyr.elf", "md5 dd56e8fe", False),
    ("Test harness", "/tmp/hfm-verify/ (scripts + prebuilts)", "two_qemu_smoke.py", False),
    ("Launcher", "launch_openbmc.sh / launch_hfm.sh", "cold-boot auto-discovery", False),
]
y = 1.8
colx = [0.7, 3.7, 9.9]
colw = [3.0, 6.2, 2.7]
for r, (a, b, c, hdr) in enumerate(rows):
    fill = ACCENT if hdr else (STEEL if r % 2 else NAVY)
    box(s, 0.7, y, 11.95, 0.72, fill=fill, line=STEEL, line_w=0.5)
    tcol = NAVY if hdr else WHITE
    text(s, colx[0]+0.15, y+0.1, colw[0]-0.2, 0.5, [[R(a, 13.5, tcol, hdr)]], anchor=MSO_ANCHOR.MIDDLE)
    fnt = "Consolas" if not hdr else "Calibri"
    text(s, colx[1]+0.15, y+0.1, colw[1]-0.2, 0.5, [[R(b, 12.5, tcol, hdr, False, fnt)]], anchor=MSO_ANCHOR.MIDDLE)
    ccol = GREEN if (not hdr) else NAVY
    text(s, colx[2]+0.15, y+0.1, colw[2]-0.2, 0.5, [[R(c, 12.5, ccol, hdr, False, "Consolas" if not hdr else "Calibri")]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.75
text(s, 0.7, 6.65, 12, 0.5,
     [[R("QEMU writes the .mtd back in place — restore the pristine image from .mtd.gz before every run.", 13, AMBER, italic=True)]])
footer(s, 6)

# ============================================================ Slide 7 — Terminal layout
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "LIVE DEMO · SETUP", "Three SSH terminals — watch each boot yourself")
cols = [
    ("Terminal 1", AMBER, "OpenBMC QEMU", "socket server · start first",
     ["cd /tmp/hfm-verify/scripts", "./launch_openbmc.sh", "# press Enter, watch", "# AST2600 boot -> login:"]),
    ("Terminal 2", ACCENT, "Zephyr QEMU", "socket client",
     ["cd /tmp/hfm-verify/scripts", "./launch_hfm.sh", "# press Enter, watch", "# boot + reverse probe"]),
    ("Terminal 3", GREEN, "Drive PLDM", "observe auto-route + forward",
     ["cd /tmp/hfm-verify/scripts", "./drive_bmc_pldm.sh", "# auto-discovered route", "# forward pldmtool"]),
]
x = 0.7
for name, col, role, sub, lines in cols:
    box(s, x, 1.75, 3.95, 4.35, fill=STEEL, line=col, line_w=1.25)
    box(s, x, 1.75, 3.95, 0.62, fill=col)
    text(s, x+0.15, 1.82, 3.65, 0.5, [[R(name, 17, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.2, 2.5, 3.6, 0.4, [[R(role, 16, WHITE, True)]])
    text(s, x+0.2, 2.9, 3.6, 0.35, [[R(sub, 12, GREY, italic=True)]])
    box(s, x+0.2, 3.35, 3.55, 2.55, fill=CODEBG, line=None)
    text(s, x+0.35, 3.5, 3.3, 2.3, [[C(ln, 12.5, ACCENT if ln.startswith("./") else (GREY if ln.startswith("#") else WHITE))] for ln in lines], line_spacing=1.25)
    x += 4.03
text(s, 0.7, 6.4, 12, 0.6,
     [[R("Order matters: T1 (server) \u2192 T2 (client) \u2192 T3 (after BMC shows the login banner). All three SSH to the same dev host.", 13, AMBER, italic=True)]])
footer(s, 7)

# ============================================================ Slide 8 — Terminals 1 & 2 boot
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "LIVE DEMO · STEP 1-2", "Boot both QEMUs — the part you can now SEE")
box(s, 0.7, 1.7, 5.85, 4.5, fill=STEEL, line=AMBER, line_w=1.0)
text(s, 0.95, 1.82, 5.4, 0.4, [[R("Terminal 1 — OpenBMC (AST2600)", 15, AMBER, True)]])
box(s, 0.95, 2.3, 5.35, 1.5, fill=CODEBG)
text(s, 1.1, 2.42, 5.05, 1.3, [
    [C("./launch_openbmc.sh", 13, ACCENT)],
    [C("<Enter>", 12.5, GREY)],
    [C("...U-Boot / kernel / systemd...", 12.5, WHITE)],
    [C("evb-ast2600 login:", 13, GREEN)],
], line_spacing=1.15)
text(s, 0.95, 3.95, 5.4, 2.1, [
    [R("• Full AST2600 boot is visible (was hidden by the one-shot script).", 13, WHITE)],
    [R("• 2nd UART (ttyS0) is exported as the MCTP socket server.", 13, WHITE)],
    [R("• Console login root / 0penBmc (T3 uses SSH, no manual login).", 13, WHITE)],
], line_spacing=1.12)
box(s, 6.75, 1.7, 5.9, 4.5, fill=STEEL, line=ACCENT, line_w=1.0)
text(s, 7.0, 1.82, 5.4, 0.4, [[R("Terminal 2 — Zephyr (RISC-V)", 15, ACCENT, True)]])
box(s, 7.0, 2.3, 5.4, 1.7, fill=CODEBG)
text(s, 7.15, 2.42, 5.1, 1.5, [
    [C("./launch_hfm.sh   <Enter>", 13, ACCENT)],
    [C("*** Booting Zephyr OS v4.3.0 ***", 12.5, WHITE)],
    [C("serial_bridge: EID 18 ... ready", 12.5, WHITE)],
    [C("BMC GetTID try 1 rc=-116, retrying", 12, AMBER)],
    [C("  ^ retries until mctpd auto-discovers", 11.5, GREY)],
], line_spacing=1.12)
text(s, 7.0, 4.15, 5.4, 2.0, [
    [R("• Zephyr = EID 18: PLDM Type 0/2 + MCTP control responder.", 13, WHITE)],
    [R("• Its requester thread probes the BMC and retries (-116) —", 13, WHITE)],
    [R("  normal until mctpd finishes auto-discovery (SetupEndpoint).", 13, GREY)],
    [R("• Then the reverse probe completes on its own — watch here.", 13, WHITE)],
], line_spacing=1.12)
footer(s, 8)

# ============================================================ Slide 9 — Demo Step 3 (forward results)
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "LIVE DEMO · STEP 3 (Terminal 3)", "Forward path — BMC discovers & polls Zephyr")
box(s, 0.7, 1.7, 11.95, 2.9, fill=CODEBG, line=GREEN, line_w=1.0)
text(s, 0.95, 1.82, 11.5, 2.7, [
    [C("BMC# mctp route          # installed by mctpd, not by hand", 14.5, ACCENT)],
    [C("        eid min 18 max 18 net 1 dev mctpserial0 mtu 68", 14, GREEN)],
    [C("BMC# pldmtool base GetTID -m 18", 14.5, ACCENT)],
    [C("        GetTID rc=0            \u2192  Response: 1", 14, GREEN)],
    [C("BMC# pldmtool base GetPLDMTypes -m 18", 14.5, ACCENT)],
    [C("        GetPLDMTypes rc=0      \u2192  base(0) + platform(2)", 14, GREEN)],
    [C("BMC# pldmtool platform GetPDR -m 18 -d 0", 14.5, ACCENT)],
    [C("        Terminus Locator PDR, recordHandle 1, TID 1, EID 18", 14, GREEN)],
    [C("BMC# pldmtool platform GetSensorReading -m 18 -i 1 --rearm 0", 14.5, ACCENT)],
    [C("        presentReading 31, Sensor Enabled, Sensor Normal", 14, GREEN)],
], line_spacing=1.05)
text(s, 0.7, 4.8, 12, 0.5, [[R("mctpd auto-installed the EID-18 route (no manual route add); drive_bmc_pldm.sh then polls over the REAL kernel AF_MCTP stack.", 13.5, WHITE)]])
box(s, 0.7, 5.4, 11.95, 1.1, fill=STEEL)
text(s, 0.95, 5.5, 11.5, 0.95, [
    [R("Point out: ", 14, AMBER, True), R("GetPDR returns a real Terminus Locator PDR and the sensor reading (die-temp = 31\u00b0C) — this is Hardware Fault Management telemetry actually flowing over the wire.", 14, WHITE)],
], line_spacing=1.12)
footer(s, 9)

# ============================================================ Slide 10 — Demo Step D (reverse / 0009)
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "LIVE DEMO · STEP 4 (watch Terminal 2)", "Reverse path — Zephyr \u2192 BMC (patch 0009)")
box(s, 0.7, 1.65, 5.85, 2.35, fill=STEEL, line=RED, line_w=1.0)
text(s, 0.95, 1.75, 5.4, 0.4, [[R("The bug (before 0009)", 15, RED, True)]])
text(s, 0.95, 2.2, 5.45, 1.75, [
    [R("Zephyr requester sent with the MCTP tag-owner (TO) bit = 0.", 13, WHITE)],
    [R("Linux AF_MCTP delivers an inbound SOM frame to a bound socket only when TO = 1; TO = 0 is treated as an orphan response and dropped (-ENOENT).", 13, WHITE)],
    [R("Symptom: rx_packets grows, tx stays 0, every call -ETIMEDOUT.", 13, GREY, False, True)],
], line_spacing=1.1)
box(s, 6.75, 1.65, 5.9, 2.35, fill=STEEL, line=GREEN, line_w=1.0)
text(s, 7.0, 1.75, 5.4, 0.4, [[R("The fix (patch 0009)", 15, GREEN, True)]])
text(s, 7.0, 2.2, 5.45, 1.75, [
    [R("In pldm.c pldm_send_request_sync():", 13, WHITE)],
    [C("MCTP_MESSAGE_TO_DST  \u2192  MCTP_MESSAGE_TO_SRC", 12.5, ACCENT)],
    [R("Requester now sets TO = 1, so the kernel routes the request to the bound responder socket.", 13, WHITE)],
    [R("Merged in RASAPI (PR #1).", 13, GREEN, False, True)],
], line_spacing=1.1)
box(s, 0.7, 4.15, 11.95, 2.05, fill=CODEBG, line=GREEN, line_w=1.0)
text(s, 0.95, 4.27, 11.5, 1.9, [
    [C("Zephyr console (answered by the stock BMC pldmd, no manual responder):", 13.5, GREY)],
    [C("<inf> serial_bridge: BMC GetTID -> 0x01", 14, GREEN)],
    [C("<inf> serial_bridge: BMC GetPLDMTypes -> byte0=0x1d", 14, GREEN)],
    [C("<inf> serial_bridge: BMC GetPLDMVersion(BASE) -> 1.0.0", 14, GREEN)],
    [C("<inf> serial_bridge: Reverse-direction PLDM probe to BMC complete", 14, GREEN)],
], line_spacing=1.06)
footer(s, 10)

# ============================================================ Slide 11 — Full auto-discovery
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "FULL AUTO-DISCOVERY", "Three fixes that make mctpd self-discover at boot")
box(s, 0.7, 1.7, 11.95, 2.5, fill=STEEL, line=GREEN, line_w=1.0)
text(s, 0.95, 1.82, 11.5, 0.4, [[R("Cold boot, zero manual steps — mctpd finds EID 18 on its own", 15, GREEN, True)]])
text(s, 0.95, 2.3, 11.5, 1.85, [
    [R("1.  ", 13.5, ACCENT, True), R("Kernel (patch 0010): ", 13.5, ACCENT, True), R("6.6.92's for_each_netdev_dump() used xa_for_each_start(), whose cursor never advances at end-of-walk — so the AF_MCTP address dump (mctp_dump_addrinfo) never emits NLMSG_DONE and mctpd spins at 100% CPU. Backport upstream cfa7fa02078d.", 13.5, WHITE)],
    [R("2.  ", 13.5, ACCENT, True), R("Zephyr control responder (patch 0011): ", 13.5, ACCENT, True), R("control replies wrongly set TO = 1; the BMC kernel treats them as new requests and the mctpd physical-addressing discovery query times out. Reply with TO = 0.", 13.5, WHITE)],
    [R("3.  ", 13.5, ACCENT, True), R("BMC systemd order (patch 0006): ", 13.5, ACCENT, True), R("mctpd snapshots the link table once at startup, so the serial link must exist first. Split into mctp-local.service (Before mctpd) + mctp-setup-endpoint.service (After mctpd).", 13.5, WHITE)],
], line_spacing=1.1)
box(s, 0.7, 4.4, 11.95, 1.9, fill=STEEL)
text(s, 0.95, 4.52, 11.5, 0.4, [[R("Verified end to end on a from-recipe rebuild", 15, ACCENT, True)]])
text(s, 0.95, 5.0, 11.5, 1.3, [
    [R("\u2713  ", 14, GREEN, True), R("mctp-local \u2192 mctpd \u2192 mctp-setup-endpoint all active; SetupEndpoint succeeds, endpoints/18 published on D-Bus", 14, WHITE)],
    [R("\u2713  ", 14, GREEN, True), R("mctpd auto-installs the kernel route (eid 18 dev mctpserial0); forward pldmtool GetTID/PDR/Sensor all rc=0", 14, WHITE)],
    [R("\u2713  ", 14, GREEN, True), R("Reverse probe (patch 0009) completes on its own — answered by the stock pldmd, no manual responder", 14, WHITE)],
], line_spacing=1.16)
footer(s, 11)

# ============================================================ Slide 12 — Recap / commands
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "RECAP", "The whole demo — three terminals")
cmds = [
    ("T1:  ./launch_openbmc.sh", "boot OpenBMC (server) \u2192 login:"),
    ("T2:  ./launch_hfm.sh", "boot Zephyr (client), watch retries"),
    ("T3:  ./drive_bmc_pldm.sh", "observe auto-route + forward pldmtool"),
    ("T3 output:  GetTID/PDR/Sensor rc=0", "forward path verified"),
    ("T2 output:  'Reverse-direction ... complete'", "reverse path (patch 0009)"),
]
y = 1.85
for i, (cmd, note) in enumerate(cmds, 1):
    box(s, 0.7, y, 0.55, 0.7, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    text(s, 0.7, y+0.08, 0.55, 0.5, [[R(str(i), 18, NAVY, True)]], align=PP_ALIGN.CENTER)
    box(s, 1.45, y, 8.2, 0.7, fill=CODEBG, line=STEEL, line_w=0.5)
    text(s, 1.6, y+0.08, 7.95, 0.55, [[C(cmd, 13.5, ACCENT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.85, y+0.08, 3.3, 0.55, [[R(note, 12.5, GREY)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.86
text(s, 0.7, 6.5, 12, 0.6,
     [[R("Full presenter script with expected output & fallbacks: ", 14, WHITE), R("DEMO_GUIDE.md", 14, ACCENT, True, False, "Consolas")]])
footer(s, 12)

# ============================================================ Slide 13 — Thanks
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 3.1, 13.333, 0.06, fill=ACCENT)
text(s, 0.9, 2.2, 11.9, 1.0, [[R("Questions & discussion", 38, WHITE, True)]])
text(s, 0.9, 3.35, 11.9, 0.6, [[R("Hardware Fault Management RAS API — MCTP / PLDM over QEMU", 18, GREY)]])
text(s, 0.9, 4.2, 11.9, 1.4, [
    [R("Repo:  OCP-China-Projects/Hardware-Fault-Management  ·  branch RASAPI", 15, WHITE)],
    [R("Design:  docs/pldm-mctp-i3c-design.md   ·   Build:  README.md   ·   Patches:  patches/README.md", 15, GREY)],
    [R("Demo host:  internal dev machine   ·   Presenter guide:  DEMO_GUIDE.md", 15, GREY)],
], line_spacing=1.2)

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "HFM_RASAPI_demo.pptx")
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
